#!/usr/bin/env python3
"""
Insert 10 demo screenshot slides into AMSI_vs_Obfuscation.pptx between
the existing 'Live Demo' overview slide (17) and the 'Summary' slide (18).
"""
import sys
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

sys.stdout.reconfigure(encoding='utf-8')

PPTX_PATH = Path("C:/VSExclude/confidence_2026/AMSI_vs_Obfuscation.pptx")
SCREEN_DIR = Path("C:/VSExclude/confidence_2026/demo_screens")
BACKUP = PPTX_PATH.with_suffix(".pptx.bak")

# Slide dimensions for this deck: 12188952 x 6858000 EMU (= 13.33" x 7.5")
SLIDE_W = Emu(12188952)
SLIDE_H = Emu(6858000)
MARGIN  = Inches(0.4)

# Colors -- match deck dark theme
BG      = RGBColor(0x0F, 0x14, 0x1A)   # dark navy
TITLE_F = RGBColor(0xE8, 0xEE, 0xF7)   # off-white
ACCENT  = RGBColor(0x3E, 0xA1, 0xFC)   # cyan blue
CAPTION = RGBColor(0xB6, 0xC1, 0xD1)   # light gray

SLIDES_TO_ADD = [
    # (filename, section, title, caption)
    ("0_amsi_bypass.PNG", "VI · Demo",
     "Layer 1 blokuje real-world AMSI bypass",
     "ETW bypass (sample 28): ramsi-com łapie 3 predykaty - PSEtwLogProvider, etwProvider, "
     "System.Management.Automation.Tracing. AMSI_RESULT_DETECTED → czerwony "
     "'ScriptContainedMaliciousContent'. Atak nie wykonuje się."),

    ("1.PNG", "VI · Demo",
     "Obfuskowany sample - ściana base64",
     "samples/malicious/14_amsi_buffer_path.ps1: cały .NET DLL (AmsiBypass.dll) zakodowany "
     "w base64, ładowany przez Reflection.Assembly::Load. Klasyczny fileless loader - "
     "payload nigdy nie ląduje na dysku."),

    ("2.PNG", "VI · Demo",
     "Layer 1 wykrywa pomimo obfuskacji",
     "Deobfuscator dekoduje base64, znajduje wewnątrz binarki strings 'AmsiScanBuffer' i "
     "'amsi.dll'. Status: AMSI BYPASS, confidence 75. Dwa niezależne predykaty: "
     "MemoryPatch (Critical) + AmsiDll (High)."),

    ("3.png", "VI · Demo",
     "Layer 1: 96% detection, 2% FP",
     "Pełen scan: 26/27 obfuskowanych bypass'ów wykrytych (real-world samples). "
     "49/50 benign installer scripts (Office, OneDrive, Autopilot) - czyste. "
     "Jeden missed: sample 27 (analizowany dalej). Jeden FP: skomplikowany installer."),

    ("4.PNG", "VI · Demo",
     "Sample 26 - radkum literal vtable hijack",
     "AMSI provider hijack: czyta CLSID-y, ładuje provider DLL-e, patchuje vtable "
     "IAntimalwareProvider (Scan slot → CloseSession). Literalne nazwy w kodzie: "
     "DllGetClassObject, WriteIntPtr, AllocHGlobal, UnmanagedFunctionPointer."),

    ("5.PNG", "VI · Demo",
     "Layer 1 łapie sample 26 - 6 predykatów",
     "ComManipulation: DllGetClassObject, GetDelegateForFunctionPointer, "
     "UnmanagedFunctionPointer. VtableManipulation: WriteIntPtr, ReadIntPtr, AllocHGlobal. "
     "Confidence 84, AMSI BYPASS. Defender by też pewnie złapał."),

    ("6.PNG", "VI · Demo",
     "Sample 27 - ewasywna wersja",
     "Robi DOKŁADNIE TO SAMO co 26, ale identyfikatory zrekonstruowane w runtime: "
     "method names z base64, 'AMSI' z char-codes 65/77/83/73, "
     "delegate attribute przez string concat. Zero literałów do match'owania."),

    ("7.png", "VI · Demo",
     "PUNCH MOMENT: Layer 1 mówi 'Clean'",
     "Deobfuscator wykrył: nic. Confidence 0, zero indicators, is_amsi_bypass: false. "
     "Statyczna analiza nie potrafi rozszyfrować że [int][char]$k[0] -eq 65 to "
     "porównanie ze znakiem 'A'. Fundamentalny limit bez symbolic execution."),

    ("8_a.PNG", "VI · Demo",
     "Layer 2 łapie sample 27 - block w konsoli ofiary",
     "Kernel widzi że proces enumeruje rejestr AMSI providers. sysmon-um: "
     "NtSuspendProcess → ocena (.ps1 z user-path = SUSPECT) → wstrzykuje czerwony "
     "tekst do konsoli ofiary → TerminateProcess. Wygląda jak natywny AMSI block."),

    ("8_b.PNG", "VI · Demo",
     "Layer 2 - pełen audit trail w sysmon-um",
     "ProcessCreate → AMSI-RECON (registry read na \\AMSI\\Providers) → SUSPENDED "
     "(NtSuspendProcess) → VERDICT: SUSPECT → TERMINATED. Niezależny kernel-level "
     "observer odporny na user-mode bypass (AMSI/ETW disabled by attacker)."),
]


def add_demo_slide(prs, screen_path, section, title, caption):
    """Add a single slide with section badge, title, image, caption."""
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    # Use Blank layout (index 6 usually) -- match existing deck which uses 'Blank'
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            blank = layout
            break
    slide = prs.slides.add_slide(blank)

    # Set black background to match deck theme
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False

    # Section badge (top-left)
    sect_tb = slide.shapes.add_textbox(MARGIN, Inches(0.25), Inches(2.5), Inches(0.4))
    sect_tf = sect_tb.text_frame
    sect_tf.margin_top = 0
    sect_tf.margin_bottom = 0
    sect_tf.margin_left = 0
    sect_tf.margin_right = 0
    p = sect_tf.paragraphs[0]
    p.text = section
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    # Title
    title_tb = slide.shapes.add_textbox(
        MARGIN, Inches(0.65), SLIDE_W - 2 * MARGIN, Inches(0.7)
    )
    title_tf = title_tb.text_frame
    title_tf.margin_top = 0
    title_tf.margin_bottom = 0
    title_tf.margin_left = 0
    title_tf.margin_right = 0
    p = title_tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_F

    # Image -- fit into remaining area
    avail_top = Inches(1.45)
    avail_h   = Inches(4.55)  # leave room for caption
    avail_w   = SLIDE_W - 2 * MARGIN

    with Image.open(screen_path) as img:
        iw, ih = img.size
    # Scale to fit
    scale_w = avail_w / Emu(int(iw * 9525))   # 9525 EMU per pixel at 96 dpi
    scale_h = avail_h / Emu(int(ih * 9525))
    scale = min(scale_w, scale_h)
    pic_w = Emu(int(iw * 9525 * scale))
    pic_h = Emu(int(ih * 9525 * scale))
    pic_x = (SLIDE_W - pic_w) // 2
    pic_y = avail_top + (avail_h - pic_h) // 2

    pic = slide.shapes.add_picture(
        str(screen_path), pic_x, pic_y, width=pic_w, height=pic_h
    )

    # Caption at bottom
    cap_tb = slide.shapes.add_textbox(
        MARGIN, Inches(6.15), SLIDE_W - 2 * MARGIN, Inches(1.2)
    )
    cap_tf = cap_tb.text_frame
    cap_tf.word_wrap = True
    cap_tf.margin_top = 0
    cap_tf.margin_bottom = 0
    p = cap_tf.paragraphs[0]
    p.text = caption
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = CAPTION

    return slide


def reorder_slides_after(prs, insert_after_idx, count):
    """Move the last `count` slides to be after position `insert_after_idx` (0-based)."""
    # python-pptx exposes sldIdLst as xml — manipulate directly
    sld_id_lst = prs.slides._sldIdLst  # internal
    sld_ids = list(sld_id_lst)
    # The newly added slides are at the end
    to_move = sld_ids[-count:]
    # Remove them from their current position
    for s in to_move:
        sld_id_lst.remove(s)
    # Insert them at the correct position
    target_position = insert_after_idx + 1
    # We need to insert at target_position in the underlying XML
    # Easiest: append after the slide at insert_after_idx
    anchor = sld_ids[insert_after_idx]
    for s in reversed(to_move):
        anchor.addnext(s)


def main():
    if not PPTX_PATH.exists():
        sys.exit(f"Not found: {PPTX_PATH}")
    if not SCREEN_DIR.exists():
        sys.exit(f"Not found: {SCREEN_DIR}")

    # Backup
    if not BACKUP.exists():
        BACKUP.write_bytes(PPTX_PATH.read_bytes())
        print(f"backup -> {BACKUP}")

    prs = Presentation(str(PPTX_PATH))
    before_count = len(prs.slides)
    print(f"loaded: {before_count} slides")

    # Add demo slides at end
    for fname, section, title, caption in SLIDES_TO_ADD:
        screen_path = SCREEN_DIR / fname
        if not screen_path.exists():
            print(f"  SKIP missing: {fname}")
            continue
        add_demo_slide(prs, screen_path, section, title, caption)
        print(f"  added: {fname} -> {title}")

    after_count = len(prs.slides)
    new_count = after_count - before_count
    print(f"added {new_count} slides")

    # Move newly added slides to be right after slide 17 (Live Demo overview)
    # Slide 17 is index 16 (0-based)
    INSERT_AFTER_IDX = 16  # i.e. after the 17th slide
    reorder_slides_after(prs, INSERT_AFTER_IDX, new_count)
    print(f"reordered: new slides inserted after slide 17")

    prs.save(str(PPTX_PATH))
    print(f"saved: {PPTX_PATH}")
    print(f"total slides: {after_count}")


if __name__ == "__main__":
    main()
