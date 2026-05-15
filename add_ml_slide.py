#!/usr/bin/env python3
"""Add an ML slide between slide 13 (Static pipeline) and current slide 14 (Layer 2 Behavioral)."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
sys.stdout.reconfigure(encoding='utf-8')

PPTX = Path("C:/VSExclude/confidence_2026/AMSI_vs_Obfuscation.pptx")
BACKUP = PPTX.with_suffix(".pptx.bak3")

SLIDE_W = Emu(12188952)
SLIDE_H = Emu(6858000)
MARGIN  = Inches(0.5)

BG       = RGBColor(0x0F, 0x14, 0x1A)
TITLE_F  = RGBColor(0xE8, 0xEE, 0xF7)
ACCENT   = RGBColor(0x3E, 0xA1, 0xFC)
GREEN    = RGBColor(0x4E, 0xD6, 0x9A)
YELLOW   = RGBColor(0xF5, 0xC8, 0x4C)
BODY     = RGBColor(0xC8, 0xD1, 0xDD)
DIM      = RGBColor(0x7A, 0x86, 0x99)

if not BACKUP.exists():
    BACKUP.write_bytes(PPTX.read_bytes())
    print(f"backup -> {BACKUP}")

prs = Presentation(str(PPTX))

# Find blank layout
blank = next((l for l in prs.slide_layouts if l.name == "Blank"), prs.slide_layouts[0])
slide = prs.slides.add_slide(blank)

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid(); bg.fill.fore_color.rgb = BG
bg.line.fill.background(); bg.shadow.inherit = False

def add_text(left, top, width, height, text, size, color, bold=False, align=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.margin_left = 0; tf.margin_right = 0
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
        if align is not None:
            p.alignment = align
    return tb

# Section badge
add_text(MARGIN, Inches(0.25), Inches(2), Inches(0.4),
         "V · Wykrywanie", 14, ACCENT, bold=True)

# Title
add_text(MARGIN, Inches(0.65), SLIDE_W - 2*MARGIN, Inches(0.7),
         "Layer 1 — Online ML (early stage)", 28, TITLE_F, bold=True)

# Subtitle
add_text(MARGIN, Inches(1.35), SLIDE_W - 2*MARGIN, Inches(0.5),
         "Heuristic rules teach the ML. ML weight grows with sample count.",
         16, DIM)

# Architecture block (left column)
col_w = Inches(5.8)
col_left = MARGIN
col_top = Inches(2.0)

add_text(col_left, col_top, col_w, Inches(0.45),
         "Architecture", 18, ACCENT, bold=True)
arch = (
    "1.  Each scan extracts 20-feature vector\n"
    "    (entropy, base64 density, reflection APIs,\n"
    "    memory APIs, line stats, ...)\n"
    "\n"
    "2.  Logistic regression scores 0.0 – 1.0\n"
    "    using current model coefficients.\n"
    "\n"
    "3.  Sample counter (persisted on disk) tracks\n"
    "    total scans accumulated.\n"
    "\n"
    "4.  Blending:\n"
    "    weight    = min(samples / 10 000, 1.0)\n"
    "    confidence = rules · (1 – w) + ml · w\n"
    "\n"
    "5.  Heuristic label saved with features\n"
    "    → future re-training on real data."
)
add_text(col_left, col_top + Inches(0.5), col_w, Inches(4.5),
         arch, 13, BODY)

# Right column: example JSON + ramp curve description
right_left = MARGIN + col_w + Inches(0.4)
right_w = SLIDE_W - 2*MARGIN - col_w - Inches(0.4)

add_text(right_left, col_top, right_w, Inches(0.45),
         "Example output (sample 26)", 18, ACCENT, bold=True)
example = (
    '"ml": {\n'
    '  "score":         0.99,\n'
    '  "weight":        0.008,    ← 79 / 10 000\n'
    '  "sample_count":  79,\n'
    '  "interpretation":\n'
    '    "early-stage model agrees with\n'
    '     rules but weight is only 0.8%",\n'
    '  "top_features": [\n'
    '    { amsi_string_count : +3.2 },\n'
    '    { entropy           : +1.9 },\n'
    '    { base64_count      : +1.4 },\n'
    '    { reflection_api    : +0.9 }\n'
    '  ]\n'
    '}'
)
add_text(right_left, col_top + Inches(0.5), right_w, Inches(4.5),
         example, 12, GREEN)

# Bottom strip: confidence ramp
add_text(MARGIN, Inches(6.15), SLIDE_W - 2*MARGIN, Inches(0.5),
         "Confidence ramp:    1 sample → 0.01 %    ·    100 → 1 %    ·    1 000 → 10 %    ·    10 000 → 100 %",
         14, YELLOW, bold=True, align=None)

# Bottom note
add_text(MARGIN, Inches(6.65), SLIDE_W - 2*MARGIN, Inches(0.6),
         "Honest demo state: rules carry > 99 % of verdict today. ML observes and learns. "
         "Architecture supports future-proof retraining; weight grows automatically as data accrues.",
         12, DIM)

# Reorder: this new slide is at end; move it after slide 13 (index 12)
sld_ids = list(prs.slides._sldIdLst)
new_slide = sld_ids[-1]
prs.slides._sldIdLst.remove(new_slide)
sld_ids[12].addnext(new_slide)
print(f"inserted new ML slide after slide 13 (index 12)")

prs.save(str(PPTX))
print(f"saved -> {PPTX}")
print(f"total slides: {len(list(prs.slides._sldIdLst))}")
