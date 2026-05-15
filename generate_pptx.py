"""
Generates AMSI_vs_Obfuscation.pptx -- ps-parser-cli edition.

Replaces the previous deck which featured PsParser.dll (C# NativeAOT).
The narrative is now pure-Rust Layer 1 (ps-parser-cli) + kernel Layer 2
(sysmon-rs). Sample 27 punch is "Suspicious not Clean" -- the best
static engine recovers `DllGetClassObject` from base64 but cannot reach
a conclusive AMSI BYPASS verdict.

Run:  python generate_pptx.py
Output: AMSI_vs_Obfuscation.pptx
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout.reconfigure(encoding="utf-8")

# ── Color palette ────────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0D, 0x1B, 0x2A)  # navy
BG_CARD     = RGBColor(0x16, 0x2A, 0x3E)
ACCENT_RED  = RGBColor(0xE6, 0x3A, 0x2F)
ACCENT_ORG  = RGBColor(0xF5, 0x8C, 0x1A)
ACCENT_GRN  = RGBColor(0x55, 0xB8, 0x6E)
ACCENT_BLU  = RGBColor(0x3E, 0xA1, 0xFC)
TEXT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GREY   = RGBColor(0xAA, 0xBB, 0xCC)
TEXT_DIM    = RGBColor(0x77, 0x88, 0x99)

SCREENS_DIR = Path("demo_screens")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=TEXT_WHITE,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = wrap
    # Replace pre-baked newlines with paragraph splits
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.text_frame.paragraphs[0] if i == 0 else tf.text_frame.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tf


def rect(slide, l, t, w, h, fill_color, line_color=None, line_w=1):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def header(slide, number, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.0, ACCENT_RED)
    txbox(slide, str(number), 0.2, 0.05, 1.5, 0.9, size=36, bold=True, color=TEXT_WHITE)
    txbox(slide, title, 1.4, 0.08, 11.5, 0.55, size=26, bold=True, color=TEXT_WHITE)
    if subtitle:
        txbox(slide, subtitle, 1.4, 0.6, 11.5, 0.35, size=13, color=RGBColor(0xFF, 0xCC, 0x99))


def bullet(slide, items, l, t, w, line_h=0.5, size=16, color=TEXT_WHITE,
           bullet_char="▸", bullet_color=ACCENT_ORG):
    """Render bullet items vertically starting at (l, t)."""
    for i, item in enumerate(items):
        y = t + i * line_h
        txbox(slide, bullet_char, l, y, 0.4, line_h, size=size, bold=True, color=bullet_color)
        txbox(slide, item, l + 0.4, y, w - 0.4, line_h, size=size, color=color)


def card(slide, l, t, w, h, title, items, title_color=ACCENT_ORG):
    rect(slide, l, t, w, h, BG_CARD)
    txbox(slide, title, l + 0.3, t + 0.15, w - 0.6, 0.45,
          size=16, bold=True, color=title_color)
    for i, line in enumerate(items):
        txbox(slide, "• " + line, l + 0.3, t + 0.7 + i * 0.4, w - 0.6, 0.4,
              size=12, color=TEXT_WHITE)


def screen(slide, name, l, t, w, h, border=True):
    path = SCREENS_DIR / name
    if not path.exists():
        # Placeholder if screen is missing
        rect(slide, l, t, w, h, BG_CARD, ACCENT_RED, 2)
        txbox(slide, f"[missing: {name}]", l, t + h / 2 - 0.2, w, 0.4,
              size=12, color=ACCENT_RED, align=PP_ALIGN.CENTER)
        return
    pic = slide.shapes.add_picture(str(path), Inches(l), Inches(t),
                                   width=Inches(w), height=Inches(h))
    if border:
        pic.line.color.rgb = TEXT_DIM
        pic.line.width = Pt(0.5)


def footer(slide, text):
    txbox(slide, text, 0, 7.05, 13.33, 0.35, size=10,
          color=TEXT_DIM, align=PP_ALIGN.CENTER, italic=True)


# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 1 — Title                                                           ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 2.8, 13.33, 1.9, ACCENT_RED)
txbox(s, "AMSI vs. Obfuscation", 0, 2.95, 13.33, 1.0,
      size=52, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
txbox(s, "A Cat-and-Mouse Game", 0, 3.85, 13.33, 0.6,
      size=26, color=TEXT_WHITE, align=PP_ALIGN.CENTER, italic=True)
txbox(s, "Radosław Kumorek  •  Kaseya", 0, 5.4, 13.33, 0.5,
      size=18, color=TEXT_GREY, align=PP_ALIGN.CENTER)
txbox(s, "Two-layer defense against PowerShell-based AMSI bypass\n"
         "Static deobfuscation (Rust) + kernel behavioral telemetry",
      0, 5.95, 13.33, 0.9, size=15, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 2 — Agenda                                                          ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "•", "Agenda", subtitle="~45 minutes")
items = [
    ("I. Script-based threats", "Why PowerShell? Stats, real-world droppers.", "2 min"),
    ("II. AMSI in 8 minutes",   "What it is, how it scans, where providers plug in.",  "8 min"),
    ("III. Bypass + obfuscation","Common techniques, AmsiProviderScanDisruption.",     "12 min"),
    ("IV. Detection — two layers","Static deobfuscation + kernel behavioral.",         "10 min"),
    ("V. Live demo",             "5 acts on real bypass samples.",                     "8 min"),
    ("VI. Q&A",                  "",                                                   "3 min"),
]
for i, (title, desc, t_hint) in enumerate(items):
    y = 1.5 + i * 0.85
    rect(s, 0.8, y, 0.15, 0.7, ACCENT_ORG)
    txbox(s, title, 1.2, y, 6.5, 0.4, size=18, bold=True, color=TEXT_WHITE)
    txbox(s, desc, 1.2, y + 0.38, 9, 0.32, size=12, color=TEXT_GREY)
    txbox(s, t_hint, 10.5, y, 2.5, 0.4, size=14, color=ACCENT_ORG, align=PP_ALIGN.RIGHT)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 3 — Why scripts are dangerous                                       ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "I", "Why scripts are dangerous")
txbox(s, "Skrypty ≠ exploity. Są jeszcze gorsze.", 0.8, 1.2, 12, 0.5,
      size=22, bold=True, color=ACCENT_ORG)

card(s, 0.7, 2.0, 4.0, 4.6, "Bliskość systemu", [
    "PowerShell, WMI, JScript",
    "Pełen dostęp do .NET",
    "Reflection bez kompilacji",
    "Living off the land",
])
card(s, 4.9, 2.0, 4.0, 4.6, "Brak wymagań", [
    "Nie wymaga kompilacji",
    "Brak pliku PE na dysku",
    "Fileless / in-memory",
    "Łatwy do dostarczenia",
])
card(s, 9.1, 2.0, 3.5, 4.6, "Trudność w detekcji", [
    "Zaufany interpreter",
    "Code = data",
    "Obfuskacja trywialna",
    "Statyka ma pułap",
])

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 4 — PowerShell + Living off the Land                                ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "I.2", "PowerShell — primary vector")
txbox(s,
      "PowerShell jest preinstalowany na każdym Windowsie od ery Win 7 SP1.\n"
      "Dla atakującego: zaufany binary, podpisany przez Microsoft, z pełnym .NET runtime.",
      0.8, 1.2, 12, 1.0, size=15, color=TEXT_GREY)

txbox(s, "Typowe role:", 0.8, 2.4, 12, 0.5, size=18, bold=True, color=ACCENT_ORG)
bullet(s, [
    "Dropper / loader   —  pobiera payload, ładuje go w pamięci",
    "Post-exploit       —  lateral movement, persistence, recon",
    "Credential dumper  —  Invoke-Mimikatz, LSASS access",
    "Ransomware stager  —  encrypt + exfiltrate + ransom note",
], 1.2, 3.0, 11, line_h=0.55, size=15)

rect(s, 0.8, 5.5, 11.7, 1.5, BG_CARD)
txbox(s, "Przykład — klasyczny dropper:", 1.0, 5.6, 11, 0.4,
      size=12, color=ACCENT_BLU, italic=True)
txbox(s, "powershell -nop -w hidden -enc <base64>",
      1.0, 6.0, 11, 0.5, size=18, bold=True, color=TEXT_WHITE)
txbox(s, "-nop = NoProfile,  -w hidden = no console window,  -enc = base64-encoded command",
      1.0, 6.5, 11, 0.4, size=11, color=TEXT_DIM, italic=True)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 5 — AMSI: czym jest                                                 ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "II", "AMSI — Anti-Malware Scan Interface")
txbox(s,
      "Microsoft-owy interfejs zaproponowany w Windows 10 (2015).\n"
      "Pozwala AV-om skanować *zawartość* skryptu PRZED jego wykonaniem.",
      0.8, 1.2, 12, 1.0, size=16, color=TEXT_WHITE)

txbox(s, "Co AMSI skanuje:", 0.8, 2.5, 12, 0.5, size=18, bold=True, color=ACCENT_ORG)
bullet(s, [
    "PowerShell — całe bloki skryptów, dynamiczne `Invoke-Expression`",
    "JScript / VBScript — Windows Script Host",
    "WMI — operacje na obiektach COM",
    "Office macros, .NET Assembly.Load()",
    "Każda aplikacja może wywołać `AmsiScanBuffer` na własnych buforach",
], 1.2, 3.1, 11, line_h=0.55, size=15)

rect(s, 0.8, 6.0, 11.7, 1.0, BG_CARD)
txbox(s,
      "Provider model — każdy AV rejestruje swoje COM DLL w HKLM\\…\\AMSI\\Providers.\n"
      "AMSI woła wszystkich providerów, jeden 'malicious' = blokada.",
      1.0, 6.15, 11.3, 0.85, size=13, color=TEXT_GREY)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 6 — AMSI flow                                                       ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "II.2", "AMSI flow")

# Flow diagram, single column centered
boxes = [
    ("PowerShell.exe",         "  user invokes a script",          ACCENT_BLU),
    ("AMSI.dll",               "  AmsiScanBuffer(content)",        ACCENT_ORG),
    ("Registered providers",   "  enumerated from registry",       ACCENT_GRN),
    ("Provider COM DLLs",      "  each scans, returns verdict",    ACCENT_GRN),
    ("Aggregated verdict",     "  any 'malicious' = block",        ACCENT_ORG),
    ("PowerShell decision",    "  CLEAN → execute   /   MAL → red error",  ACCENT_RED),
]
for i, (title, sub, c) in enumerate(boxes):
    y = 1.4 + i * 0.85
    rect(s, 3.5, y, 6.3, 0.7, BG_CARD, c, 2)
    txbox(s, title, 3.7, y + 0.05, 6, 0.35, size=15, bold=True, color=c)
    txbox(s, sub, 3.7, y + 0.4, 6, 0.3, size=11, color=TEXT_GREY)
    if i < len(boxes) - 1:
        txbox(s, "▼", 6.3, y + 0.7, 0.7, 0.2, size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 7 — Bypass techniques taxonomy                                      ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "III", "AMSI bypass — taxonomy")

card(s, 0.5, 1.4, 4.0, 5.7, "Patch w pamięci", [
    "VirtualProtect + WriteProcessMemory",
    "Patchuje AmsiScanBuffer entry",
    "Returns AMSI_RESULT_CLEAN",
    "→ Matt Graeber (2016)",
    "→ Tal Liberman variants",
    "Pierwszy znany bypass; nadal działa na niezaaktualizowanych targetach",
])
card(s, 4.7, 1.4, 4.0, 5.7, "Reflection bypass", [
    "[Ref].Assembly.GetType('AmsiUtils')",
    ".GetField('amsiInitFailed')",
    ".SetValue($null, $true)",
    "Wyłącza AMSI session flag",
    "→ One-liner, najprostszy",
    "Patchowany przez MS w 2019 ale stale wraca w wariantach",
])
card(s, 8.9, 1.4, 4.0, 5.7, "COM/vtable hijack", [
    "DllGetClassObject(provider CLSID)",
    "GetDelegateForFunctionPointer",
    "WriteIntPtr → patch vtable",
    "Każdy COM-side hook",
    "→ AmsiProviderScanDisruption",
    "Lokalny, in-process, omija nawet poprawne AMSI provider'y",
])

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 8 — AmsiProviderScanDisruption                                      ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "III.2", "AmsiProviderScanDisruption  —  case study")
txbox(s,
      "github.com/radkum/AmsiProviderScanDisruption  —  technika opublikowana 5 lat temu.",
      0.8, 1.15, 12, 0.4, size=14, color=ACCENT_BLU, italic=True)

txbox(s, "Idea:", 0.8, 1.7, 12, 0.4, size=18, bold=True, color=ACCENT_ORG)
bullet(s, [
    "Wylistuj registered providerów z HKLM\\…\\AMSI\\Providers",
    "Załaduj ich COM DLL przez `DllGetClassObject`",
    "Wyciągnij IAntimalwareProvider COM object",
    "Patchuj vtable: `IAntimalwareProvider::Scan` → return CLEAN",
    "Każdy kolejny AMSI scan wraca CLEAN — bez modyfikacji amsi.dll",
], 1.2, 2.2, 11, line_h=0.5, size=14)

rect(s, 0.8, 5.4, 11.7, 1.6, BG_CARD, ACCENT_RED, 2)
txbox(s, "Co to znaczy dla obrońcy:", 1.0, 5.5, 11, 0.4,
      size=13, color=ACCENT_RED, italic=True, bold=True)
txbox(s,
      "Atakujący nie tknął amsi.dll. Wszystkie EDR-y patrzące na patch amsi.dll widzą czysto.\n"
      "Provider DLL też nietknięty — patch jest w vtable, w pamięci procesu.\n"
      "Statyczna sygnatura złapie literalny kod. Ewasywna wersja (sample 27) — nie.",
      1.0, 5.85, 11.5, 1.1, size=12, color=TEXT_WHITE)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 9 — Obfuscation techniques                                          ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "III.3", "Obfuscation — utrudnianie detekcji")

obf = [
    ("Base64 encoding",      "powershell -enc <blob>  /  FromBase64String"),
    ("String concatenation", "'Amsi'+'Init'+'Failed' — naiwna ale skuteczna"),
    ("-f format operator",   "\"{2}{0}{1}\" -f 'Init','Failed','Amsi' → 'AmsiInitFailed'"),
    ("[char] casts / arrays","[char[]]@(65,109,115,105) -join ''  →  'Amsi'"),
    ("Custom class method",  "class C { static [string] X() { 'am' } } ;  [C]::X()+'siInit'+'Failed'  →  'amsiInitFailed'"),
    ("XOR + byte array",     "@(0x62,0x6e,0x70,0x6a) | %{[char]($_ -bxor 0x03)}"),
    ("Backtick obfuscation", "Inv`oke-Ex`pression  — backticki ignorowane w identifierach"),
    ("Reverse string slice", "'tpircSekovnI'[-1..-12] -join ''"),
    ("Env-var indexing",     "$env:ComSpec[14,15,16,17] -join ''  →  custom"),
    ("Embedded .NET DLL",    "[Reflection.Assembly]::Load(FromBase64String(...))  — fileless"),
    ("Runtime-built names",  "char-code compare zamiast literałów — niedeobfuskowalne statycznie"),
]
for i, (name, ex) in enumerate(obf):
    y = 1.3 + i * 0.55
    rect(s, 0.5, y, 4.3, 0.45, BG_CARD)
    txbox(s, name, 0.7, y + 0.05, 4, 0.35, size=14, bold=True, color=ACCENT_ORG)
    txbox(s, ex, 5.0, y + 0.05, 8.0, 0.45, size=12, color=TEXT_WHITE)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 10 — Two-layer defense overview                                     ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "IV", "Two-layer defense")
txbox(s, "Statystyka spotyka behaviour. Każda warstwa łapie co druga przepuści.",
      0.8, 1.15, 12, 0.4, size=15, color=ACCENT_ORG, italic=True)

# Layer 1
rect(s, 0.5, 1.8, 6.1, 5.2, BG_CARD, ACCENT_BLU, 2)
txbox(s, "Layer 1  —  Inline AMSI provider", 0.7, 1.9, 5.7, 0.5,
      size=17, bold=True, color=ACCENT_BLU)
txbox(s, "PowerShell → AMSI → ramsi-com.dll (Rust COM)\n"
         "                   → ps-parser-cli engine",
      0.7, 2.4, 5.7, 0.95, size=12, color=TEXT_WHITE)
bullet(s, [
    "Statyczna deobfuskacja",
    "Format operator, char arrays",
    "Rekursywny base64 (8 layers)",
    "30+ predykatów AMSI bypass",
    "Blokuje BEFORE execution",
    "Werdykt: Clean / Suspicious / BYPASS",
], 0.7, 3.5, 5.7, line_h=0.45, size=13)
txbox(s, "Detection: 27/28 = 96 % static", 0.7, 6.4, 5.7, 0.4,
      size=14, bold=True, color=ACCENT_GRN)

# Layer 2
rect(s, 6.8, 1.8, 6.1, 5.2, BG_CARD, ACCENT_RED, 2)
txbox(s, "Layer 2  —  Kernel behavioral", 7.0, 1.9, 5.7, 0.5,
      size=17, bold=True, color=ACCENT_RED)
txbox(s, "Kernel driver (sysmon-rs, no_std)\n"
         "→ userspace daemon (sysmon-um)",
      7.0, 2.4, 5.7, 0.95, size=12, color=TEXT_WHITE)
bullet(s, [
    "PsSetCreateProcessNotifyEx",
    "PsSetLoadImageNotify",
    "CmRegisterCallback (registry)",
    "Filter \\AMSI\\ in registry path",
    "Suspend → evaluate → terminate",
    "Werdykt na podstawie behaviour",
], 7.0, 3.5, 5.7, line_h=0.45, size=13)
txbox(s, "Catches what static cannot", 7.0, 6.4, 5.7, 0.4,
      size=14, bold=True, color=ACCENT_GRN)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 11 — Layer 1 details                                                ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "IV.1", "Layer 1  —  Static deobfuscation (pure Rust)")

# Architecture diagram
txbox(s, "Architecture", 0.8, 1.2, 6, 0.4, size=18, bold=True, color=ACCENT_ORG)
boxes = [
    ("PowerShell process",   ACCENT_BLU),
    ("amsi.dll  ::  AmsiScanBuffer",  ACCENT_ORG),
    ("ramsi-com.dll  (Rust COM provider)", ACCENT_GRN),
    ("ps-parser-cli engine  (in-process)", ACCENT_GRN),
    ("AMSI verdict  →  block / allow",   ACCENT_RED),
]
for i, (title, c) in enumerate(boxes):
    y = 1.7 + i * 0.75
    rect(s, 0.8, y, 5.7, 0.55, BG_CARD, c, 2)
    txbox(s, title, 1.0, y + 0.07, 5.5, 0.4, size=13, bold=True, color=c)
    if i < len(boxes) - 1:
        txbox(s, "▼", 0.8, y + 0.55, 5.7, 0.18, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Right column — what ps-parser-cli does
txbox(s, "ps-parser-cli internals", 7.0, 1.2, 6, 0.4, size=18, bold=True, color=ACCENT_ORG)
bullet(s, [
    "pest grammar — recovers AST + evaluates",
    "String methods: Replace, ToUpper, …",
    "Format operator -f, [char] casts",
    "Recursive base64 (up to 8 layers)",
    "UTF-16LE projection for .NET DLLs",
    "30+ predykatów: amsi.dll, AmsiScanBuffer,",
    "  amsiInitFailed, AmsiUtils, ETW, WLDP, …",
    "Combo bonuses + scoring algorithm",
    "Output: JSON  →  ramsi-com",
], 7.0, 1.7, 6, line_h=0.5, size=13)

txbox(s, "4.7 MB statyczna binarka  •  zero .NET dependencies",
      7.0, 6.4, 6, 0.4, size=13, italic=True, color=ACCENT_BLU)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 12 — ps-parser safe-evaluation                                      ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "IV.1b", "ps-parser  —  safe evaluation")
txbox(s,
      "Pure data transforms execute.  Side-effectful calls return typed errors — never silently skipped.\n"
      "API celowo nazwane `safe_eval` — granica zaufania jest wbudowana w architekturę.",
      0.8, 1.1, 12, 0.7, size=13, color=ACCENT_ORG, italic=True)

# LEFT  —  Evaluated
rect(s, 0.5, 1.95, 6.1, 4.7, BG_CARD, ACCENT_GRN, 2)
txbox(s, "✓ Evaluated", 0.7, 2.05, 5.7, 0.4, size=17, bold=True, color=ACCENT_GRN)
txbox(s, "whitelist — pure, deterministic, no I/O", 0.7, 2.45, 5.7, 0.3,
      size=10, color=TEXT_GREY, italic=True)
left_items = [
    ("String methods",    "Replace, Substring, Insert, ToUpper/Lower, Trim, Split, PadLeft/Right"),
    ("Static helpers",    "[Convert]::FromBase64String, [Text.Encoding]::UTF8.GetString"),
    ("Operators",         "+, -, *, /, %, -bxor, -band, -bor, -shl, -shr, -replace, -split"),
    ("Format operator",   "\"{0}{1}\" -f 'foo','bar'  →  'foobar'"),
    ("Variables + scope", "$script:, $global:, $local:, $env:  (via Variables::env())"),
    ("PS classes",        "class C { static [string] X() { 'am' } }  →  [C]::X()  ⇒  'am'"),
]
for i, (name, ex) in enumerate(left_items):
    y = 2.85 + i * 0.62
    txbox(s, name, 0.7, y, 5.6, 0.3, size=12, bold=True, color=TEXT_WHITE)
    txbox(s, ex, 0.7, y + 0.3, 5.6, 0.3, size=9, color=TEXT_GREY)

# RIGHT  —  Skipped (typed error)
rect(s, 6.7, 1.95, 6.1, 4.7, BG_CARD, ACCENT_RED, 2)
txbox(s, "✗ Skipped", 6.9, 2.05, 5.7, 0.4, size=17, bold=True, color=ACCENT_RED)
txbox(s, "side-effects → MethodError / CommandError", 6.9, 2.45, 5.7, 0.3,
      size=10, color=TEXT_GREY, italic=True)
right_items = [
    ("File I/O",          "Get-Content, Set-Content, Out-File, [IO.File]::*"),
    ("Network",           "Invoke-WebRequest, Net.WebClient::DownloadString"),
    ("Process spawn",     "Start-Process, Invoke-Expression, iex"),
    ("Dynamic code",      "Add-Type, [Reflection.Assembly]::Load(...)"),
    ("Registry",          "Get-ItemProperty, [Microsoft.Win32.Registry]::*"),
    ("Filesystem cmdlets","Get-ChildItem, New-Item, Remove-Item, dowolna mutacja stanu"),
]
for i, (name, ex) in enumerate(right_items):
    y = 2.85 + i * 0.62
    txbox(s, name, 6.9, y, 5.6, 0.3, size=12, bold=True, color=TEXT_WHITE)
    txbox(s, ex, 6.9, y + 0.3, 5.6, 0.3, size=9, color=TEXT_GREY)

# Punchline at bottom
rect(s, 0.5, 6.75, 12.3, 0.55, BG_CARD)
txbox(s,
      "Skutek:  deobfuscator widzi WSZYSTKO co da się policzyć z czystych operacji.  "
      "Atak który zależy od I/O lub spawning nie wykonuje się — i jednocześnie jest sygnałem.",
      0.6, 6.83, 12.1, 0.4, size=11, italic=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 13 — Layer 2 details                                                ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "IV.2", "Layer 2  —  Kernel behavioral telemetry")

txbox(s, "Kernel driver (sysmon-rs, no_std)", 0.8, 1.2, 6, 0.4,
      size=18, bold=True, color=ACCENT_ORG)
bullet(s, [
    "PsSetCreateProcessNotifyRoutineEx",
    "    → capture cmdline at process spawn",
    "PsSetLoadImageNotifyRoutine",
    "    → flag amsi.dll / wldp.dll loads",
    "CmRegisterCallbackEx",
    "    → filter Registry access by path",
    "Sends events via \\\\.\\SysMon to userspace",
], 0.8, 1.7, 6.0, line_h=0.45, size=13)

txbox(s, "Userspace daemon (sysmon-um)", 7.0, 1.2, 6, 0.4,
      size=18, bold=True, color=ACCENT_ORG)
bullet(s, [
    "Polling \\\\.\\SysMon at 200 ms",
    "Tracks pid → cmdline mapping",
    "On AMSI registry recon:",
    "    1.  NtSuspendProcess(pid)",
    "    2.  call ps-parser-cli on cmdline",
    "    3.  Suspicious → log,  BYPASS → kill",
    "Audit log per scanned process",
], 7.0, 1.7, 6.0, line_h=0.45, size=13)

rect(s, 0.8, 5.7, 11.7, 1.4, BG_CARD, ACCENT_RED, 2)
txbox(s, "Co Layer 2 widzi a Layer 1 nie:", 1.0, 5.8, 11, 0.4,
      size=14, bold=True, color=ACCENT_RED)
txbox(s, "Każde otwarcie klucza rejestru ze ścieżką \\AMSI\\ — niezależnie od tego "
         "jak nazwa klucza była zbudowana w kodzie.\n"
         "Behaviour pozostaje, nawet gdy literalne identyfikatory są rekonstruowane w runtime.",
      1.0, 6.15, 11.5, 0.95, size=12, color=TEXT_WHITE)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 13 — Demo intro                                                     ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 2.5, 13.33, 2.5, ACCENT_ORG)
txbox(s, "LIVE DEMO", 0, 2.7, 13.33, 0.95, size=58, bold=True,
      color=TEXT_WHITE, align=PP_ALIGN.CENTER)
txbox(s, "5 actów — od literalnego bypass'a po behavioral catch",
      0, 3.7, 13.33, 0.6, size=22, color=TEXT_WHITE,
      align=PP_ALIGN.CENTER, italic=True)
txbox(s, "Act 1: Layer 1 blocks an obfuscated bypass",   0, 5.3, 13.33, 0.4,
      size=15, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Act 2: Detection rate sweep — 27/28 = 96 %",   0, 5.7, 13.33, 0.4,
      size=15, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Act 3: Literal radkum technique caught live",  0, 6.1, 13.33, 0.4,
      size=15, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Act 4: Evasive variant — only Suspicious",     0, 6.5, 13.33, 0.4,
      size=15, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Act 5: Behavioral catches the residual",       0, 6.9, 13.33, 0.4,
      size=15, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 14 — Act 1: obfuscated source                                       ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "V.1", "Act 1  —  Obfuscated AMSI bypass (sample 14)")
txbox(s, "Cały skompilowany .NET DLL zakodowany w base64, ładowany przez Reflection.Assembly::Load.",
      0.8, 1.15, 12, 0.4, size=14, color=ACCENT_ORG, italic=True)
screen(s, "1_obfuscated_sample.png", 1.5, 1.7, 10.3, 4.8)
footer(s, "samples\\malicious\\14_amsi_buffer_path.ps1  —  first 8 lines.   "
          "Line 3 = entire AmsiBypass.dll as base64.")

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 15 — Act 1 catch                                                    ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "V.1b", "Act 1  —  ps-parser-cli verdict")
txbox(s, "Rekurencyjny base64 deobfuscator wyciąga literały z UTF-16LE #US section.",
      0.8, 1.15, 12, 0.4, size=14, color=ACCENT_ORG, italic=True)
screen(s, "2_pscli_sample14_detection.png", 2.5, 1.7, 8.0, 4.5)

rect(s, 0.8, 6.3, 11.7, 0.75, BG_CARD, ACCENT_GRN, 2)
txbox(s, "status = AMSI BYPASS    •    confidence = 100    •    indicators: "
         "AmsiDll, MemoryPatch, VtableManipulation, TelemetryFunction",
      0.9, 6.4, 11.5, 0.55, size=13, bold=True, color=ACCENT_GRN, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 16 — Act 2: corpus sweep                                            ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "V.2", "Act 2  —  Detection rate sweep")
txbox(s, "Tym samym silnikiem na pełnym zestawie 28 obfuskowanych sample'ów.",
      0.8, 1.15, 12, 0.4, size=14, color=ACCENT_ORG, italic=True)
screen(s, "3_pscli_full_corpus.png", 0.8, 1.7, 11.7, 5.0)

rect(s, 0.8, 6.85, 11.7, 0.45, BG_CARD)
txbox(s, "27 × AMSI BYPASS    •    1 × Suspicious  (sample 27 — patrz Akt 4)    •    0 × Clean",
      0.9, 6.9, 11.5, 0.35, size=13, bold=True, color=ACCENT_GRN, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 17 — Act 3 source                                                   ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "V.3", "Act 3  —  Literal radkum technique (sample 26)")
txbox(s, "Adaptacja github.com/radkum/AmsiProviderScanDisruption — bez obfuskacji.",
      0.8, 1.15, 12, 0.4, size=14, color=ACCENT_ORG, italic=True)
screen(s, "4_sample26_source.png", 1.5, 1.7, 10.3, 4.7)
footer(s, "Literalne DllGetClassObject, Marshal.WriteIntPtr, AllocHGlobal — wszystko na wierzchu.")

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 18 — Act 3 live block                                               ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "V.3b", "Act 3  —  Live block in PowerShell")
txbox(s, "Provider zwraca AMSI_RESULT_DETECTED. PS odmawia wykonania.",
      0.8, 1.15, 12, 0.4, size=14, color=ACCENT_ORG, italic=True)
screen(s, "5_amsi_red_error.png", 1.5, 1.7, 10.3, 4.6)

rect(s, 0.8, 6.4, 11.7, 0.7, BG_CARD, ACCENT_RED, 2)
txbox(s, "Standardowy AMSI block message — nie nasza nakładka.  "
         "Skrypt sprzed 5 lat → blokowany dziś przez własny detector autora.",
      0.9, 6.5, 11.5, 0.5, size=12, italic=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 19 — Act 4 source (sample 27)                                       ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "V.4", "Act 4  —  Evasive variant (sample 27)")
txbox(s, "To samo zachowanie co sample 26 — każdy identyfikator zrekonstruowany w runtime.",
      0.8, 1.15, 12, 0.4, size=14, color=ACCENT_ORG, italic=True)
screen(s, "6_sample27_source.png", 1.0, 1.7, 11.3, 4.4)

bullet(s, [
    "Type name z char-code compare:  [int][char]$k[0] -eq 65   (= 'A')",
    "Method name z base64:  'RGxsR2V0Q2xhc3NPYmplY3Q='  →  DllGetClassObject",
    "Attribute name z konkatenacji:  'Unmanaged' + 'FunctionPointer'",
], 0.8, 6.2, 12, line_h=0.35, size=12)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 20 — Act 4 PUNCH: ps-parser-cli Suspicious                          ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "V.4b", "Act 4  —  ⭐ PUNCH  ⭐  ps-parser-cli says 'Suspicious'")
txbox(s, "Best static engine recovers DllGetClassObject from base64. But cannot reach BYPASS verdict.",
      0.8, 1.15, 12, 0.4, size=14, color=ACCENT_ORG, italic=True)
screen(s, "7_pscli_sample27_suspicious.png", 3.0, 1.7, 7.3, 4.6)

rect(s, 0.8, 6.4, 11.7, 0.75, BG_CARD, ACCENT_ORG, 2)
txbox(s, "status = Suspicious    •    confidence = 24    •    is_amsi_bypass = false    •    "
         "indicators: ComManipulation × 2",
      0.9, 6.5, 11.5, 0.55, size=13, bold=True, color=ACCENT_ORG, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 21 — Act 4c: Layer 1 doesn't block                                  ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "V.4c", "Act 4  —  Suspicious is NOT a block")
txbox(s, "Provider returns CLEAN to AMSI (is_amsi_bypass=false). PowerShell executes the script.",
      0.8, 1.15, 12, 0.4, size=14, color=ACCENT_ORG, italic=True)
screen(s, "7_pscli_sample27_not_blocked.png", 1.5, 1.7, 10.3, 4.5)

rect(s, 0.8, 6.3, 11.7, 0.85, BG_CARD, ACCENT_RED, 2)
txbox(s, "To jest pułap statyki w działaniu.  Silnik widzi pęknięcie, nie ma dowodu na bypass.",
      0.9, 6.4, 11.5, 0.35, size=13, bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
txbox(s, "Behavioral observation closes the gap.",
      0.9, 6.75, 11.5, 0.35, size=13, italic=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 22 — Act 5: Layer 2 catches                                         ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "V.5", "Act 5  —  Layer 2 (kernel) catches the recon")
txbox(s, "Każdy OpenSubKey z \\AMSI\\ w ścieżce  →  event AMSI-RECON w sysmon-um.",
      0.8, 1.15, 12, 0.4, size=14, color=ACCENT_ORG, italic=True)
screen(s, "8_sample27_layer2_caught.png", 1.0, 1.7, 11.3, 4.7)

rect(s, 0.8, 6.55, 11.7, 0.55, BG_CARD, ACCENT_GRN, 2)
txbox(s, "Żadna legitymna aplikacja nie czyta tych kluczy.  Decyzja w 10 ms.  Definitywna.",
      0.9, 6.65, 11.5, 0.4, size=13, bold=True, color=ACCENT_GRN, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 23 — Summary: the gradient                                          ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
header(s, "VI", "Summary — detection gradient on sample 27")

# Three columns, but middle column is the punch
col_y = 1.5
col_h = 5.0
col_w = 4.0
gap = 0.3

# Left: AST-only parsers (na potrzeby narracji "AST-only parser" - hypothetical)
rect(s, 0.5, col_y, col_w, col_h, BG_CARD, TEXT_DIM, 2)
txbox(s, "Naive static parser", 0.7, col_y + 0.2, col_w - 0.4, 0.4,
      size=15, bold=True, color=TEXT_DIM)
txbox(s, "AST + literal-string match", 0.7, col_y + 0.65, col_w - 0.4, 0.3,
      size=11, color=TEXT_DIM, italic=True)
bullet(s, [
    "Widzi: literal AMSI",
    "Widzi: literal AmsiUtils",
    "Widzi: amsiInitFailed",
], 0.7, col_y + 1.2, col_w - 0.4, line_h=0.45, size=12, color=TEXT_GREY)
rect(s, 0.7, col_y + col_h - 1.1, col_w - 0.4, 0.85, BG_DARK, TEXT_DIM, 1)
txbox(s, "Clean", 0.7, col_y + col_h - 1.05, col_w - 0.4, 0.5,
      size=24, bold=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)
txbox(s, "no signal", 0.7, col_y + col_h - 0.55, col_w - 0.4, 0.3,
      size=10, italic=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Middle: ps-parser-cli
rect(s, 0.5 + col_w + gap, col_y, col_w, col_h, BG_CARD, ACCENT_ORG, 2)
txbox(s, "ps-parser-cli", 0.7 + col_w + gap, col_y + 0.2, col_w - 0.4, 0.4,
      size=15, bold=True, color=ACCENT_ORG)
txbox(s, "Rust + recursive base64 + 30 predykatów", 0.7 + col_w + gap, col_y + 0.65, col_w - 0.4, 0.3,
      size=11, color=ACCENT_ORG, italic=True)
bullet(s, [
    "Wyciąga DllGetClassObject z b64",
    "Łapie ComManipulation",
    "Sygnał — ale bez literalnego",
    "AMSI w żadnej formie",
], 0.7 + col_w + gap, col_y + 1.2, col_w - 0.4, line_h=0.45, size=12)
rect(s, 0.7 + col_w + gap, col_y + col_h - 1.1, col_w - 0.4, 0.85, BG_DARK, ACCENT_ORG, 1)
txbox(s, "Suspicious", 0.7 + col_w + gap, col_y + col_h - 1.05, col_w - 0.4, 0.5,
      size=22, bold=True, color=ACCENT_ORG, align=PP_ALIGN.CENTER)
txbox(s, "widzi pęknięcie — ale nie dowód", 0.7 + col_w + gap, col_y + col_h - 0.55, col_w - 0.4, 0.3,
      size=10, italic=True, color=ACCENT_ORG, align=PP_ALIGN.CENTER)

# Right: Layer 2
rect(s, 0.5 + 2 * (col_w + gap), col_y, col_w, col_h, BG_CARD, ACCENT_GRN, 2)
txbox(s, "Layer 2 (kernel)", 0.7 + 2 * (col_w + gap), col_y + 0.2, col_w - 0.4, 0.4,
      size=15, bold=True, color=ACCENT_GRN)
txbox(s, "Behavioral — runtime observation", 0.7 + 2 * (col_w + gap), col_y + 0.65, col_w - 0.4, 0.3,
      size=11, color=ACCENT_GRN, italic=True)
bullet(s, [
    "Process otwiera \\AMSI\\ key",
    "Niezależne od literałów",
    "Niezależne od obfuskacji",
    "Widzi realny dostęp",
], 0.7 + 2 * (col_w + gap), col_y + 1.2, col_w - 0.4, line_h=0.45, size=12)
rect(s, 0.7 + 2 * (col_w + gap), col_y + col_h - 1.1, col_w - 0.4, 0.85, BG_DARK, ACCENT_GRN, 1)
txbox(s, "AMSI BYPASS", 0.7 + 2 * (col_w + gap), col_y + col_h - 1.05, col_w - 0.4, 0.5,
      size=20, bold=True, color=ACCENT_GRN, align=PP_ALIGN.CENTER)
txbox(s, "definitive, 10 ms", 0.7 + 2 * (col_w + gap), col_y + col_h - 0.55, col_w - 0.4, 0.3,
      size=10, italic=True, color=ACCENT_GRN, align=PP_ALIGN.CENTER)

txbox(s,
      "Statyka ma fundamentalny pułap.  Każda warstwa ma swoje miejsce.\n"
      "Defense in depth nie jest opcjonalne — jest konieczne.",
      0.5, 6.65, 12.3, 0.7, size=14, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

# ╔════════════════════════════════════════════════════════════════════════════╗
# ║  Slide 24 — Q&A / References                                               ║
# ╚════════════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(BLANK)
add_bg(s)
rect(s, 0, 1.2, 13.33, 1.5, ACCENT_RED)
txbox(s, "Q & A", 0, 1.35, 13.33, 1.2, size=72, bold=True,
      color=TEXT_WHITE, align=PP_ALIGN.CENTER)

txbox(s, "Resources", 0.8, 3.3, 12, 0.5, size=22, bold=True, color=ACCENT_ORG)
bullet(s, [
    "ps-parser   —  crates.io/crates/ps-parser     (Rust PowerShell evaluator)",
    "AmsiProviderScanDisruption   —  github.com/radkum/AmsiProviderScanDisruption",
    "windows-kernel-rs   —  github.com/radkum/windows-kernel-rs",
    "AMSI documentation   —  learn.microsoft.com/en-us/windows/win32/amsi/",
    "Confidence   —  this project, full source available on request",
], 1.0, 3.9, 11.5, line_h=0.5, size=14)

txbox(s, "Radosław Kumorek   •   Kaseya   •   radoslaw.kumorek@kaseya.com",
      0, 6.8, 13.33, 0.4, size=13, color=TEXT_GREY, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = Path("AMSI_vs_Obfuscation.pptx")
# Backup existing
if out.exists():
    bak = out.with_suffix(".pptx.bak_pre_pscli")
    if not bak.exists():
        bak.write_bytes(out.read_bytes())
        print(f"backup -> {bak}")
prs.save(str(out))
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
