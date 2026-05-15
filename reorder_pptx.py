#!/usr/bin/env python3
"""
Reorder/edit slides:
1. Swap slides 9 and 10 (so general bypass overview comes BEFORE specific vtable hijack)
2. Delete slide 14 (ML deep-dive) -- we don't have ML layer; just mention briefly elsewhere
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn
sys.stdout.reconfigure(encoding='utf-8')

PPTX = Path("C:/VSExclude/confidence_2026/AMSI_vs_Obfuscation.pptx")
BACKUP = PPTX.with_suffix(".pptx.bak2")

if not BACKUP.exists():
    BACKUP.write_bytes(PPTX.read_bytes())
    print(f"backup -> {BACKUP}")

prs = Presentation(str(PPTX))

# Reorder via sldIdLst
sld_id_lst = prs.slides._sldIdLst
sld_ids = list(sld_id_lst)
print(f"before: {len(sld_ids)} slides")

# 1. Swap index 8 and 9 (slides 9 and 10 in 1-based)
s9 = sld_ids[8]
s10 = sld_ids[9]
# Move s10 right after s9's original predecessor
# Cleanest: remove s9, then s10 stays in its position (now at index 8),
# then re-insert s9 after s10
sld_id_lst.remove(s9)
s10.addnext(s9)
print("swapped slides 9 and 10")

# 2. Delete slide 14 (after swap, slide 14 is still at index 13 -- not affected by 9/10 swap)
# Refresh list after swap
sld_ids = list(sld_id_lst)
s14 = sld_ids[13]
sld_id_lst.remove(s14)
print("deleted slide 14 (ML deep-dive)")

# Note: the relationship to the slide part is still in the package; python-pptx
# doesn't delete the underlying part. The slide just won't be in the presentation
# anymore. PowerPoint cleans this up on save / reopen.

prs.save(str(PPTX))
print(f"saved -> {PPTX}")
print(f"total slides now: {len(list(prs.slides._sldIdLst))}")
